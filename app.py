import streamlit as st
import pandas as pd
from io import BytesIO

from utils import storage, data_processing as dp
from utils.diagram import render_formation_svg

st.set_page_config(page_title="OffenseIQ Scout Cards", layout="wide")

PLAYER_LABELS = ["X", "Z", "Y", "H", "F", "T", "TE", "RB", "QB", "C", "G", "T2", "TE2", "Slot", "Motion"]


def get_field_options(master_df, col):
    if master_df.empty or col not in master_df.columns:
        return []
    vals = sorted(v for v in master_df[col].dropna().astype(str).unique() if v.strip())
    return vals


def main():
    st.title("🏈 OffenseIQ — Offensive Scout Cards")
    st.caption(
        "Upload Hudl playlists, teach it how each formation is drawn once, "
        "set how your defense aligns to it, and it builds scout cards from there on out."
    )

    master_df = storage.load_master_plays()
    library = storage.load_formation_library()
    rules = storage.load_rules()

    tab_upload, tab_library, tab_rules, tab_cards = st.tabs(
        ["📤 Upload Playlists", "✏️ Formation Library", "🛡️ Alignment Rules", "📋 Scout Cards"]
    )

    # ---------------------------------------------------------- UPLOAD TAB
    with tab_upload:
        st.subheader("Upload one or more Hudl playlist exports")
        uploads = st.file_uploader(
            "Playlist .xlsx files", type=["xlsx"], accept_multiple_files=True
        )
        if uploads:
            if st.button("Process & add to master database", type="primary"):
                total_added = 0
                for f in uploads:
                    try:
                        df = dp.read_playlist(f)
                    except Exception as e:
                        st.error(f"Couldn't read {f.name}: {e}")
                        continue
                    master_df = storage.append_plays(df, source_label=f.name)
                    total_added += len(df)
                st.success(f"Added {total_added} offensive plays from {len(uploads)} file(s).")
                st.rerun()

        st.divider()
        if master_df.empty:
            st.info("No plays in the master database yet. Upload a playlist above to get started.")
        else:
            st.subheader("What's in the database so far")
            st.write(f"**{len(master_df)}** total offensive plays across **{master_df['__source_file'].nunique()}** uploaded file(s).")
            summary = dp.formation_summary(master_df)
            summary["Drawn?"] = summary["OFF FORM"].apply(lambda f: "✅" if f in library else "—")
            st.dataframe(summary, use_container_width=True, hide_index=True)
            undrawn = summary[summary["Drawn?"] == "—"]
            if not undrawn.empty:
                st.warning(
                    f"{len(undrawn)} formation(s) haven't been drawn yet — head to **Formation Library** to teach them: "
                    + ", ".join(undrawn["OFF FORM"].tolist())
                )

    # ---------------------------------------------------------- LIBRARY TAB
    with tab_library:
        st.subheader("Teach OffenseIQ how a formation lines up")
        st.caption("Do this once per formation name — every future upload with that same formation name reuses it automatically.")

        formation_names = sorted(set(dp.formation_summary(master_df)["OFF FORM"].tolist()) | set(library.keys())) if not master_df.empty else sorted(library.keys())
        new_name = st.text_input("...or type a new formation name not yet seen in an upload")
        if new_name.strip():
            formation_names = sorted(set(formation_names) | {new_name.strip()})

        if not formation_names:
            st.info("Upload a playlist first, or type a formation name above, to start drawing.")
        else:
            selected = st.selectbox("Formation", formation_names, key="lib_select")

            # copy-from helper for near-duplicate formations (e.g. TRIPS GS vs TRIPS GW)
            copy_from_options = ["— none —"] + [f for f in library.keys() if f != selected]
            copy_from = st.selectbox("Start from an existing formation's diagram (optional)", copy_from_options)

            state_key = f"editor_points_{selected}"
            if state_key not in st.session_state:
                if selected in library:
                    st.session_state[state_key] = list(library[selected]["points"])
                else:
                    st.session_state[state_key] = []

            if copy_from != "— none —" and st.button(f"Copy points from '{copy_from}'"):
                st.session_state[state_key] = [dict(p) for p in library[copy_from]["points"]]
                st.rerun()

            col_editor, col_preview = st.columns([1, 1])

            with col_editor:
                st.markdown("**Add a position**")
                c1, c2, c3 = st.columns([1, 1, 1])
                with c1:
                    label = st.selectbox("Label", PLAYER_LABELS, key=f"label_{selected}")
                with c2:
                    x = st.slider("Left ↔ Right", 0, 100, 50, key=f"x_{selected}")
                with c3:
                    y = st.slider("Deep ↔ Backfield", 0, 100, 60, key=f"y_{selected}")
                if st.button("➕ Add position", key=f"add_{selected}"):
                    st.session_state[state_key].append({"label": label, "x": x, "y": y})
                    st.rerun()

                st.markdown("**Current positions**")
                if st.session_state[state_key]:
                    for i, p in enumerate(st.session_state[state_key]):
                        pc1, pc2 = st.columns([4, 1])
                        pc1.write(f"{p['label']}  (x={p['x']}, y={p['y']})")
                        if pc2.button("🗑️", key=f"del_{selected}_{i}"):
                            st.session_state[state_key].pop(i)
                            st.rerun()
                else:
                    st.caption("No positions added yet.")

                notes = st.text_area("Notes (splits, alerts, etc.)", value=library.get(selected, {}).get("notes", ""), key=f"notes_{selected}")

                if st.button("💾 Save formation", type="primary", key=f"save_{selected}"):
                    storage.save_formation(selected, st.session_state[state_key], notes)
                    st.success(f"Saved '{selected}'. It'll auto-load on every future scout card.")

            with col_preview:
                st.markdown("**Live preview**")
                st.markdown(render_formation_svg(st.session_state[state_key], title=selected), unsafe_allow_html=True)

    # ---------------------------------------------------------- RULES TAB
    with tab_rules:
        st.subheader("Default alignment rules")
        st.caption("These apply automatically based on formation family, strength, hash, etc. — most-specific matching rule wins.")

        field_choices = {
            "FORM FAMILY": ["Any"] + get_field_options(master_df, "FORM FAMILY"),
            "OFF STR": ["Any"] + get_field_options(master_df, "OFF STR"),
            "FIELD/BOUNDARY": ["Any"] + get_field_options(master_df, "FIELD/BOUNDARY"),
            "ST/ WK": ["Any"] + get_field_options(master_df, "ST/ WK"),
            "BACKFIELD": ["Any"] + get_field_options(master_df, "BACKFIELD"),
        }

        with st.form("new_default_rule"):
            st.markdown("**When...**")
            m1, m2, m3, m4, m5 = st.columns(5)
            match = {
                "FORM FAMILY": m1.selectbox("Form Family", field_choices["FORM FAMILY"]),
                "OFF STR": m2.selectbox("Strength", field_choices["OFF STR"]),
                "FIELD/BOUNDARY": m3.selectbox("Field/Bnd", field_choices["FIELD/BOUNDARY"]),
                "ST/ WK": m4.selectbox("St/Wk", field_choices["ST/ WK"]),
                "BACKFIELD": m5.selectbox("Backfield", field_choices["BACKFIELD"]),
            }
            st.markdown("**...we call:**")
            c1, c2, c3 = st.columns(3)
            front = c1.text_input("Front call (e.g. Over Sam)")
            technique = c2.text_input("DL techniques (e.g. 4i-3-5-9)")
            coverage = c3.text_input("Coverage note (e.g. Solo)")
            submitted = st.form_submit_button("➕ Add default rule")
            if submitted:
                if not any([front, technique, coverage]):
                    st.error("Enter at least a front call, technique, or coverage note.")
                else:
                    storage.add_default_rule(match, {"front": front, "technique": technique, "coverage": coverage})
                    st.success("Rule added.")
                    st.rerun()

        st.markdown("**Existing default rules**")
        if rules["defaults"]:
            for i, rule in enumerate(rules["defaults"]):
                match_str = ", ".join(f"{k}={v}" for k, v in rule["match"].items()) or "any formation"
                call = rule["call"]
                call_str = " / ".join(v for v in [call.get("front"), call.get("technique"), call.get("coverage")] if v)
                rc1, rc2 = st.columns([5, 1])
                rc1.write(f"**If** {match_str} → **{call_str}**")
                if rc2.button("🗑️", key=f"delrule_{i}"):
                    storage.delete_default_rule(i)
                    st.rerun()
        else:
            st.caption("No default rules yet.")

        st.divider()
        st.subheader("Per-formation overrides")
        st.caption("Locks in one exact call for one specific formation name, regardless of the default rules.")
        formation_names = sorted(set(dp.formation_summary(master_df)["OFF FORM"].tolist()) | set(library.keys())) if not master_df.empty else sorted(library.keys())
        if formation_names:
            sel = st.selectbox("Formation", formation_names, key="override_select")
            existing = rules["formation_overrides"].get(sel, {})
            oc1, oc2, oc3 = st.columns(3)
            o_front = oc1.text_input("Front call", value=existing.get("front", ""), key=f"of_{sel}")
            o_tech = oc2.text_input("DL techniques", value=existing.get("technique", ""), key=f"ot_{sel}")
            o_cov = oc3.text_input("Coverage note", value=existing.get("coverage", ""), key=f"oc_{sel}")
            bc1, bc2 = st.columns(2)
            if bc1.button("💾 Save override", type="primary"):
                storage.set_formation_override(sel, {"front": o_front, "technique": o_tech, "coverage": o_cov})
                st.success(f"Saved override for {sel}.")
                st.rerun()
            if sel in rules["formation_overrides"] and bc2.button("Clear override"):
                storage.clear_formation_override(sel)
                st.rerun()
        else:
            st.info("Upload plays or add a formation in the Library tab first.")

    # ---------------------------------------------------------- SCOUT CARDS TAB
    with tab_cards:
        st.subheader("Scout cards")
        if master_df.empty:
            st.info("Upload a playlist to generate scout cards.")
        else:
            summary = dp.formation_summary(master_df)
            options = ["All formations"] + summary["OFF FORM"].tolist()
            choice = st.selectbox("Show", options)
            targets = summary["OFF FORM"].tolist() if choice == "All formations" else [choice]

            for form_name in targets:
                stats = dp.formation_tendencies(master_df, form_name)
                if not stats:
                    continue
                call, source = storage.resolve_call(form_name, stats["representative_row"], rules)

                with st.container(border=True):
                    left, right = st.columns([1, 1.3])
                    with left:
                        points = library.get(form_name, {}).get("points", [])
                        st.markdown(render_formation_svg(points, title=form_name), unsafe_allow_html=True)
                        if not points:
                            st.caption("⚠️ Not drawn yet — add it in Formation Library.")
                        notes = library.get(form_name, {}).get("notes", "")
                        if notes:
                            st.caption(f"📝 {notes}")

                    with right:
                        st.markdown(f"### {form_name}  ·  {stats['form_family']}")
                        st.write(
                            f"**{stats['total_plays']} plays seen**  |  "
                            f"Run **{stats['run_pct']}%** / Pass **{stats['pass_pct']}%**"
                            + (f"  |  Motion **{stats['motion_rate']}%**" if stats.get("motion_rate") else "")
                        )
                        if stats["avg_distance"] is not None:
                            st.write(f"Avg distance-to-go: **{stats['avg_distance']}**")
                        if stats["strength_pct"]:
                            st.write("Strength: " + ", ".join(f"{k} {v}%" for k, v in stats["strength_pct"].items()))
                        if stats["field_boundary_pct"]:
                            st.write("Field/Boundary: " + ", ".join(f"{k} {v}%" for k, v in stats["field_boundary_pct"].items()))
                        if len(stats["top_backfields"]):
                            st.write("Top backfields: " + ", ".join(f"{k} ({v})" for k, v in stats["top_backfields"].items()))
                        if len(stats["top_plays"]):
                            st.write("Top plays: " + ", ".join(f"{k} ({v})" for k, v in stats["top_plays"].items()))

                        st.markdown("---")
                        st.markdown(f"**Our call** — _{source}_")
                        qc1, qc2, qc3 = st.columns(3)
                        q_front = qc1.text_input("Front", value=(call or {}).get("front", ""), key=f"card_front_{form_name}")
                        q_tech = qc2.text_input("Technique", value=(call or {}).get("technique", ""), key=f"card_tech_{form_name}")
                        q_cov = qc3.text_input("Coverage", value=(call or {}).get("coverage", ""), key=f"card_cov_{form_name}")
                        save_default = st.checkbox("Save this as the default for this formation", key=f"card_savedef_{form_name}")
                        if st.button("Apply", key=f"card_apply_{form_name}"):
                            if save_default:
                                storage.set_formation_override(form_name, {"front": q_front, "technique": q_tech, "coverage": q_cov})
                                st.success("Saved as this formation's default call.")
                                st.rerun()
                            else:
                                st.info("Applied for this view only (not saved as a default).")

            st.divider()
            if st.button("📥 Export all scout cards to PowerPoint", type="primary"):
                pptx_bytes = build_pptx(master_df, library, rules)
                st.download_button(
                    "Download scout_cards.pptx",
                    data=pptx_bytes,
                    file_name="scout_cards.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )


def build_pptx(master_df, library, rules):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    summary = dp.formation_summary(master_df)

    for form_name in summary["OFF FORM"].tolist():
        stats = dp.formation_tendencies(master_df, form_name)
        if not stats:
            continue
        call, source = storage.resolve_call(form_name, stats["representative_row"], rules)
        slide = prs.slides.add_slide(blank)

        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.6))
        title_box.text_frame.text = f"{form_name}  ({stats['form_family']})"
        title_box.text_frame.paragraphs[0].font.size = Pt(28)
        title_box.text_frame.paragraphs[0].font.bold = True

        body = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9), Inches(5.5))
        tf = body.text_frame
        tf.word_wrap = True
        lines = [
            f"{stats['total_plays']} plays seen  |  Run {stats['run_pct']}% / Pass {stats['pass_pct']}%",
        ]
        if stats["avg_distance"] is not None:
            lines.append(f"Avg distance-to-go: {stats['avg_distance']}")
        if stats["strength_pct"]:
            lines.append("Strength: " + ", ".join(f"{k} {v}%" for k, v in stats["strength_pct"].items()))
        if stats["field_boundary_pct"]:
            lines.append("Field/Boundary: " + ", ".join(f"{k} {v}%" for k, v in stats["field_boundary_pct"].items()))
        if len(stats["top_backfields"]):
            lines.append("Top backfields: " + ", ".join(f"{k} ({v})" for k, v in stats["top_backfields"].items()))
        if len(stats["top_plays"]):
            lines.append("Top plays: " + ", ".join(f"{k} ({v})" for k, v in stats["top_plays"].items()))
        notes = library.get(form_name, {}).get("notes", "")
        if notes:
            lines.append(f"Notes: {notes}")
        lines.append("")
        lines.append(f"Our call ({source}):")
        if call:
            lines.append(
                "  ".join(
                    f"{label}: {call.get(key)}"
                    for label, key in [("Front", "front"), ("Technique", "technique"), ("Coverage", "coverage")]
                    if call.get(key)
                )
                or "Not set"
            )
        else:
            lines.append("Not set")

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


if __name__ == "__main__":
    main()
